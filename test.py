import random
import argparse
import os
import pandas as pd
import pysolr
from docx.api import Document
from pptx import Presentation


def getText(f):
    doc = Document(f)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)


def jpeg_res(filename):
    # open image for reading in binary mode
    with open(filename, 'rb') as img_file:
        # height of image (in 2 bytes) is at 164th position
        img_file.seek(163)

        # read the 2 bytes
        a = img_file.read(2)

        # calculate height
        height = (a[0] << 8) + a[1]

        # next 2 bytes is width
        a = img_file.read(2)

        # calculate width
        width = (a[0] << 8) + a[1]
        # print("The resolution of the image is",width,"x",height)
        return '{} x {}'.format(width, height)


def processFile(namafile, p):
    typeFile = os.path.splitext(namafile)[1][1:].strip()
    if typeFile in ("doc", "docx"):  # filter document file
        data = getText(namafile)
        d = dict();
        #d['nama file'] = namafile
        d['resource'] = p
        d['data'] = data
        return d

    elif typeFile in ("xls", "xlsx"):  # filter excel files
        excel_file = namafile
        xlsx = pd.ExcelFile(excel_file)
        PC_sheets = []
        for sheet in xlsx.sheet_names:
            PC_sheets.append(xlsx.parse(sheet))
        data = ' '.join([str(elem) for elem in PC_sheets])
        d = dict();
        #d['resource'] = namafile
        d['resource'] = p
        d['data'] = data
        return d

    elif typeFile in ("ppt", "pptx"):  # filter ppt
        prs = Presentation(namafile)
        data1 = []
        for slide in prs.slides:

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        data1.append(run.text)
                        data = str("".join(data1))
        d = dict();
        #d['resource'] = namafile
        d['resource'] = p
        d['data'] = data
        return d

    elif typeFile in ("jpg", "png"):  # filter extension type
        folder = os.getcwd()
        dirpath = os.path.dirname(namafile)
        # list_path1 = folder+(str(os.path.join(dirpath,f)))[1:]
        # extension1 = os.path.splitext(f)[1][1:].strip()
        data = jpeg_res(namafile)
        d = dict();
        #d['resource'] = namafile
        d['resource'] = p
        d['size'] = data
        return d


def processFolder(namafolder,p):
    curentfolder = os.getcwd()
    list_path = []
    data = []
    for path in os.listdir(namafolder):
        full_path = (curentfolder + '\\' + str(namafolder) + "\\" + path)
        if os.path.isfile(full_path):
            list_path.append(full_path)
    for filename in list_path:
        hasil = processFile(filename,p)
        if hasil is not None:
            data.append(hasil)
    return data


def push_solr(result):
    solr_url = 'http://192.168.15.205:8983/solr/kotekaman_text'
    solr = pysolr.Solr(solr_url, always_commit=True)
    solr.add(result)


def push_solr_folder(data):
    for p in data:
        solr_url = 'http://192.168.15.205:8983/solr/kotekaman_text'
        solr = pysolr.Solr(solr_url, always_commit=True)
        solr.add(p)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()

    parser.add_argument('-f', type=str, required=False,
                        help="processing file")

    parser.add_argument('-d', type=str, required=False,
                        help="processing folder")
    parser.add_argument('-p', type=str, required=False,
                        help="processing folder")

    args = parser.parse_args()

    f = args.f
    d = args.d
    p = args.p

    if f is not None and p is not None and d is None:
        result = processFile(f, p)
        #print(result)
        push_solr(result)
    else:
        result1 = processFolder(d,p)
        #print(result1)
        push_solr_folder(result1)
